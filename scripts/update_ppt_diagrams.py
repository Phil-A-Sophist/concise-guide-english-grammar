#!/usr/bin/env python3
"""Replace old diagram images in PPT slides with new table + tree PNG pairs.

Each old single-image slide becomes two slides:
  1. Slide with the labeling table PNG
  2. Slide with the tree diagram PNG

Non-image slides are preserved as-is.
"""

import json
import copy
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN

sys.stdout.reconfigure(encoding='utf-8')

PROJECT_ROOT = Path(__file__).parent.parent
DATA_FILE = PROJECT_ROOT / "data" / "static" / "ppt_diagram_data.json"
TABLE_DIR = PROJECT_ROOT / "data" / "static" / "ppt-diagrams" / "tables"
TREE_DIR = PROJECT_ROOT / "data" / "static" / "ppt-diagrams" / "trees"
PPT_DIR = Path(r"C:\Users\irphy\OneDrive - University of Colorado Colorado Springs\++Cloud\+++++Spring26\F25 PPTs")

PPT_FILES = {
    "adverbials": PPT_DIR / "F25 Class (Adverbials).pptx",
    "nominals": PPT_DIR / "F25 - Class (Nominals).pptx",
    "adjectivals": PPT_DIR / "F25 - Class (Adjectivals).pptx",
}

# Preview slide title overrides
PREVIEW_TITLES = {
    "adv_04_preview": "Preview: Adjectivals \u2014 Relative Clauses",
    "adv_05_preview": "Preview: Adjectivals \u2014 Relative Clauses",
    "adv_06_preview": "Preview: Adjectivals \u2014 Relative Clauses",
    "adj_10_preview": "Preview: Nominals \u2014 Complement Clauses",
    "adj_11_preview": "Preview: Nominals \u2014 Complement Clauses",
    "adj_12_preview": "Preview: Nominals \u2014 Complement Clauses",
}


def build_slide_map(data):
    """Build mapping: (ppt_name, slide_num) -> sentence_id."""
    slide_map = {}
    for s in data["sentences"]:
        for loc in s["slides"]:
            key = (loc["ppt"], loc["slide"])
            slide_map[key] = s["id"]
    return slide_map


def get_slide_title(slide):
    """Extract the first non-empty text from the slide."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if t and len(t) > 2:
                    return t
    return ""


def has_image(slide):
    """Check if slide has any images."""
    return any(s.shape_type == 13 for s in slide.shapes)


def get_image_shapes(slide):
    """Get all image shapes from a slide."""
    return [s for s in slide.shapes if s.shape_type == 13]


def create_image_slide(prs, layout, title_text, subtitle_text, image_path, slide_width, slide_height):
    """Create a new slide with a title and centered image."""
    slide = prs.slides.add_slide(layout)

    # Clear any placeholder text
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:  # title
            ph.text = title_text
            for run in ph.text_frame.paragraphs[0].runs:
                run.font.size = Pt(24)
                run.font.bold = True
        elif ph.placeholder_format.idx == 1:  # subtitle/body
            ph.text = subtitle_text
            for run in ph.text_frame.paragraphs[0].runs:
                run.font.size = Pt(14)
                run.font.italic = True

    # Add the image centered
    from PIL import Image
    img = Image.open(image_path)
    img_w, img_h = img.size

    # Calculate dimensions to fit within slide (with margins)
    max_w = slide_width - Inches(2)
    max_h = slide_height - Inches(2.5)  # leave room for title

    aspect = img_w / img_h
    if img_w / max_w > img_h / max_h:
        # Width is the constraint
        width = max_w
        height = int(width / aspect)
    else:
        # Height is the constraint
        height = max_h
        width = int(height * aspect)

    left = (slide_width - width) // 2
    top = Inches(1.8)

    slide.shapes.add_picture(str(image_path), left, top, width, height)
    return slide


def process_ppt(ppt_name, ppt_path, slide_map, data):
    """Process a single PPT file."""
    print(f"\n{'='*60}")
    print(f"Processing: {ppt_path.name}")
    print(f"{'='*60}")

    prs = Presentation(str(ppt_path))
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Get a blank layout (or the first layout)
    blank_layout = None
    title_layout = None
    for layout in prs.slide_layouts:
        if layout.name == 'Blank':
            blank_layout = layout
        if 'Title' in layout.name and 'Content' not in layout.name:
            title_layout = layout
    if blank_layout is None:
        blank_layout = prs.slide_layouts[-1]  # last layout as fallback

    # Collect which slides need replacement
    original_slides = list(prs.slides)
    replacements = []

    for i, slide in enumerate(original_slides):
        slide_num = i + 1
        key = (ppt_name, slide_num)
        if key in slide_map and has_image(slide):
            sid = slide_map[key]
            title = get_slide_title(slide)
            replacements.append({
                "slide_index": i,
                "slide_num": slide_num,
                "sentence_id": sid,
                "original_title": title,
            })
            print(f"  Slide {slide_num}: {title[:50]} -> {sid}")
        else:
            print(f"  Slide {slide_num}: keep as-is")

    if not replacements:
        print("  No replacements needed.")
        return

    # Strategy: Build a new presentation from scratch, copying non-image slides
    # and inserting table+tree pairs for image slides.
    # python-pptx doesn't support inserting slides at arbitrary positions easily,
    # so we'll use XML manipulation.

    # Actually, the simplest approach: for each image slide, remove the old image
    # and add the table PNG. Then duplicate the slide and replace with tree PNG.
    # But python-pptx can't easily duplicate slides either.

    # Best approach: rebuild slide order using XML manipulation.
    # For each replacement slide:
    #   1. Remove old image from the slide, add table PNG -> becomes the "table slide"
    #   2. Create a new slide after it with tree PNG -> "tree slide"

    from lxml import etree

    new_slides_to_insert = []  # (after_slide_index, sentence_id, title)

    for rep in replacements:
        idx = rep["slide_index"]
        sid = rep["sentence_id"]
        slide = prs.slides[idx]
        title = rep["original_title"]

        # Check for preview title override
        if sid in PREVIEW_TITLES:
            new_title = PREVIEW_TITLES[sid]
        else:
            new_title = title

        # Remove old images from this slide
        shapes_to_remove = get_image_shapes(slide)
        for shape in shapes_to_remove:
            sp = shape._element
            sp.getparent().remove(sp)

        # Update title if needed
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t == title and new_title != title:
                        for run in para.runs:
                            run.text = ""
                        if para.runs:
                            para.runs[0].text = new_title
                        break
                break

        # Add table PNG to this slide
        table_path = TABLE_DIR / f"{sid}.png"
        if table_path.exists():
            from PIL import Image
            img = Image.open(table_path)
            img_w, img_h = img.size
            aspect = img_w / img_h

            max_w = slide_width - Inches(1.5)
            max_h = slide_height - Inches(2.0)
            if img_w / max_w > img_h / max_h:
                width = max_w
                height = int(width / aspect)
            else:
                height = max_h
                width = int(height * aspect)

            left = (slide_width - width) // 2
            top = Inches(1.5)
            slide.shapes.add_picture(str(table_path), left, top, width, height)
            print(f"    Table added to slide {rep['slide_num']}")

        # Queue a new tree slide to insert after
        new_slides_to_insert.append({
            "after_index": idx,
            "sid": sid,
            "title": new_title,
        })

    # Now create new tree slides (add at end, then reorder via XML)
    tree_slide_map = {}  # idx -> new slide rId

    for ns in new_slides_to_insert:
        sid = ns["sid"]
        title_text = ns["title"]
        tree_path = TREE_DIR / f"{sid}.png"

        # Add a new blank slide
        new_slide = prs.slides.add_slide(blank_layout)

        # Add title text box
        # Inches, Pt already imported at top level
        txBox = new_slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(8), Inches(0.8))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = title_text
        run.font.size = Pt(24)
        run.font.bold = True

        # Add tree PNG
        if tree_path.exists():
            from PIL import Image
            img = Image.open(tree_path)
            img_w, img_h = img.size
            aspect = img_w / img_h

            max_w = slide_width - Inches(1.5)
            max_h = slide_height - Inches(2.0)
            if img_w / max_w > img_h / max_h:
                width = max_w
                height = int(width / aspect)
            else:
                height = max_h
                width = int(height * aspect)

            left = (slide_width - width) // 2
            top = Inches(1.5)
            new_slide.shapes.add_picture(str(tree_path), left, top, width, height)

        tree_slide_map[ns["after_index"]] = new_slide

    # Reorder slides: for each replacement, the tree slide should come right after its table slide
    # The new slides were added at the end. We need to move them.
    # python-pptx slide ordering is controlled by sldIdLst in presentation.xml

    PML_NS = '{http://schemas.openxmlformats.org/presentationml/2006/main}'
    sldIdLst = prs._element.find(f'{PML_NS}sldIdLst')
    if sldIdLst is None:
        sldIdLst = prs._element.find('.//sldIdLst')

    if sldIdLst is not None:
        sldIds = list(sldIdLst)
        original_count = len(sldIds) - len(new_slides_to_insert)

        # Build new order
        new_order = []
        tree_sldIds = sldIds[original_count:]  # the newly added slides
        tree_idx = 0

        for i, sldId in enumerate(sldIds[:original_count]):
            new_order.append(sldId)
            if i in tree_slide_map and tree_idx < len(tree_sldIds):
                new_order.append(tree_sldIds[tree_idx])
                tree_idx += 1

        # Any remaining tree slides
        while tree_idx < len(tree_sldIds):
            new_order.append(tree_sldIds[tree_idx])
            tree_idx += 1

        # Clear and rebuild
        for child in list(sldIdLst):
            sldIdLst.remove(child)
        for sldId in new_order:
            sldIdLst.append(sldId)

        print(f"  Reordered: {len(new_order)} slides (was {len(sldIds)})")

    # Save
    output_path = ppt_path.parent / f"{ppt_path.stem} - Updated{ppt_path.suffix}"
    prs.save(str(output_path))
    print(f"  Saved: {output_path.name}")


def main():
    with open(DATA_FILE) as f:
        data = json.load(f)

    slide_map = build_slide_map(data)

    for ppt_name, ppt_path in PPT_FILES.items():
        if not ppt_path.exists():
            print(f"SKIP: {ppt_path} not found")
            continue
        process_ppt(ppt_name, ppt_path, slide_map, data)

    print("\nDone! Updated PPTs saved with ' - Updated' suffix.")


if __name__ == "__main__":
    main()
